#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
build_daily_ppt.py (patched v3)
- Honors a top-level NEEDED list when --cols is not specified.
- Column aliasing preserved (e.g., "t_started" -> "t_start_utc"/"t_start").
- Table headers show the requested names exactly.
"""

from __future__ import annotations

import argparse
from pathlib import Path
from datetime import datetime
import re
import os

import pandas as pd
import pytz

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# === User-requested columns (used when --cols is not supplied) ===
NEEDED = ["object", "site", "detected", "t_started"]

# Default columns (fallback only if nothing requested matches at all)
DEFAULT_COLS = ["sub_id", "n", "t_start", "t_end", "span_min", "ra_deg", "dec_deg", "status"]

# Column alias map: requested key -> acceptable source column names (in order of preference)
COLUMN_ALIASES = {
    "t_started": ["t_start_utc", "t_start", "t_start_iso", "t_start_local"],
    "t_ended":   ["t_end_utc", "t_end", "t_end_iso", "t_end_local"],
    "n":         ["n_paths", "n", "png_count"],
    "object":    ["object"],
    "site":      ["site"],
    "detected":  ["detected"],
}


def sanitize(s: str) -> str:
    return re.sub(r"[^0-9A-Za-z._-]+", "_", str(s).strip()).strip("_") or "NA"


def resolve_date_label(date_arg: str, tzname: str) -> str:
    if date_arg.lower() == "today":
        tz = pytz.timezone(tzname)
        now_local = datetime.now(tz)
        return now_local.strftime("%Y-%m-%d")
    return date_arg


def load_manifest(day_out: Path) -> pd.DataFrame:
    manifest_path = day_out / "group_manifest.csv"
    if not manifest_path.is_file():
        raise FileNotFoundError(f"group_manifest.csv not found at: {manifest_path}")

    df = pd.read_csv(manifest_path)

    # Ensure required columns exist
    for c in ["object", "site", "sub_id"]:
        if c not in df.columns:
            raise KeyError(f"Required column '{c}' not found in manifest: {manifest_path}")

    # Derive gif_path if missing or not pointing to an existing file
    def derive_path(row):
        p = str(row.get("gif_path", "") or "").strip()
        if p:
            pth = Path(p)
            if not pth.is_absolute():
                pth = (day_out / p).resolve()
            if pth.is_file():
                return str(pth)

        obj = sanitize(row["object"])
        site = sanitize(row["site"])
        sub = int(row["sub_id"])
        subdir = day_out / f"{site}__{obj}"
        base = f"{obj}__{site}__sub{sub:02d}"

        cand_annoted = subdir / f"{base}_annoted.gif"
        if cand_annoted.is_file():
            return str(cand_annoted)

        cand_annotated = subdir / f"{base}_annotated.gif"
        if cand_annotated.is_file():
            return str(cand_annotated)

        cand_default = subdir / f"{base}.gif"
        return str(cand_default)

    if "gif_path" not in df.columns:
        df["gif_path"] = ""
    df["gif_path"] = df.apply(derive_path, axis=1)

    return df


def add_title_slide(prs: Presentation, title: str, subtitle: str, date_label: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title + Subtitle
    slide.shapes.title.text = f"{title} — {date_label}"
    sub = subtitle or "Automatically generated from daily_workflow outputs"
    slide.placeholders[1].text = sub


def _resolve_columns_and_headers(gdf: pd.DataFrame, requested_cols: list[str]) -> tuple[list[str], list[str]]:
    """
    Return (headers, keys):
    - headers: names to display in the table header (requested names)
    - keys:    actual DataFrame column keys to read data from (may be aliases)
    Uses partial match: missing requested columns are simply skipped.
    Fallback to DEFAULT_COLS ONLY if none of the requested columns matched at all.
    """
    headers: list[str] = []
    keys: list[str] = []

    req = requested_cols or []
    for name in req:
        # exact match
        if name in gdf.columns:
            headers.append(name)
            keys.append(name)
            continue
        # alias match
        for alias in COLUMN_ALIASES.get(name, []):
            if alias in gdf.columns:
                headers.append(name)  # display requested name
                keys.append(alias)    # read from alias column
                break

    # If nothing matched at all, fallback
    if not keys:
        fallback = [c for c in DEFAULT_COLS if c in gdf.columns]
        if not fallback and len(gdf.columns) > 0:
            fallback = list(gdf.columns[:8])
        headers = fallback[:]
        keys = fallback[:]

    return headers, keys


def add_group_slide(prs: Presentation, group_key, gdf: pd.DataFrame, cols: list[str], layout: dict):
    obj, site = group_key
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only or similar

    # Positions / sizes
    left = Inches(layout["left_margin_in"])
    top_title = Inches(layout["top_margin_in"] - 0.2)
    table_top = Inches(layout["top_margin_in"] + 0.3)
    table_width = Inches(layout["table_width_in"])
    table_height = Inches(layout["table_height_in"])

    # Title: use placeholder if available
    title_shape = slide.shapes.title
    if title_shape is not None:
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        p.text = f"{site} — {obj}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(40, 40, 40)
    else:
        tbox = slide.shapes.add_textbox(left, top_title, table_width, Inches(0.45))
        tf = tbox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{site} — {obj}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(40, 40, 40)

    # Resolve columns and headers
    headers, keys = _resolve_columns_and_headers(gdf, cols)
    sub_df = gdf[keys].copy()
    n_rows = min(len(sub_df), layout["rows_per_table"])
    rows = max(1, n_rows + 1)
    cols_n = max(1, len(keys))

    # Create table
    table = slide.shapes.add_table(rows, cols_n, left, table_top, table_width, table_height).table

    # Header row (display requested names even if aliases used)
    for j, head in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(head)
        for para in cell.text_frame.paragraphs:
            para.font.bold = True

    # Body rows
    for i in range(n_rows):
        for j, colkey in enumerate(keys):
            val = sub_df.iloc[i][colkey]
            text = "" if pd.isna(val) else str(val)
            table.cell(i + 1, j).text = text

    # GIF grid
    gifs = []
    for _, row in gdf.iterrows():
        path = str(row.get("gif_path", "")).strip()
        if path:
            pth = Path(path)
            if not pth.is_absolute():
                pth = (layout["day_out"] / path).resolve()
            if pth.is_file():
                cap = f"sub{int(row['sub_id']):02d} | n={int(row['n'])}" if "n" in row else f"sub{int(row['sub_id']):02d}"
                gifs.append((pth, cap))
        if len(gifs) >= layout["max_gifs_per_slide"]:
            break

    if gifs:
        gif_width_in = layout["gif_width_in"]
        gap_in = layout["gif_gap_in"]
        usable_width_in = layout["table_width_in"]
        cols_fit = max(1, int((usable_width_in + gap_in) // (gif_width_in + gap_in)))
        x_positions = [layout["left_margin_in"] + c * (gif_width_in + gap_in) for c in range(cols_fit)]
        y_start_in = layout["top_margin_in"] + layout["table_height_in"] + 0.6

        r = 0
        c = 0
        for gif_path, caption in gifs:
            left_pic = Inches(x_positions[c])
            top_pic = Inches(y_start_in + r * (gif_width_in + 0.55))
            slide.shapes.add_picture(str(gif_path), left_pic, top_pic, width=Inches(gif_width_in))
            cap_box = slide.shapes.add_textbox(left_pic, top_pic + Inches(gif_width_in + 0.05), Inches(gif_width_in), Inches(0.35))
            para = cap_box.text_frame.paragraphs[0]
            para.text = caption
            para.font.size = Pt(12)
            para.alignment = PP_ALIGN.CENTER
            c += 1
            if c >= cols_fit:
                c = 0
                r += 1


def main():
    ap = argparse.ArgumentParser(description="Build a PPT daily report from daily_workflow.py outputs.")
    ap.add_argument("--thumbs-root", default="./thumbs", help="Root directory where daily_workflow stores outputs.")
    ap.add_argument("--date", default="today", help='Date in YYYY-MM-DD or "today" (Australia/Melbourne by default).')
    ap.add_argument("--timezone", default="Australia/Melbourne", help="Timezone when --date=today.")
    ap.add_argument("--ppt-outdir", default=None, help="Where to save PPT (default: thumbs-root).")
    ap.add_argument("--ppt-name-template", default="{date}_daily_report.pptx", help="Output filename template.")
    ap.add_argument("--title", default="Daily Asteroid Report", help="Title for the title slide.")
    ap.add_argument("--subtitle", default="", help="Subtitle for the title slide.")
    ap.add_argument("--cols", nargs="*", default=None, help="Columns for the table on each slide. If omitted, use NEEDED.")
    ap.add_argument("--max-gifs-per-slide", type=int, default=12, help="Max GIFs per slide.")
    ap.add_argument("--rows-per-table", type=int, default=20, help="Max table rows per slide.")
    ap.add_argument("--gif-width-in", type=float, default=2.3, help="GIF width (inches).")
    ap.add_argument("--gif-gap-in", type=float, default=0.2, help="Gap between GIFs (inches).")
    ap.add_argument("--left-margin-in", type=float, default=0.5, help="Left margin (inches).")
    ap.add_argument("--top-margin-in", type=float, default=0.7, help="Top margin (inches).")
    ap.add_argument("--table-width-in", type=float, default=9.0, help="Table width (inches).")
    ap.add_argument("--table-height-in", type=float, default=2.2, help="Table height (inches).")
    ap.add_argument("--limit-groups", type=int, default=None, help="Limit number of (object,site) groups (debug).")
    args = ap.parse_args()

    date_label = resolve_date_label(args.date, args.timezone)
    thumbs_root = Path(args.thumbs_root).resolve()
    day_out = (thumbs_root / f"{date_label}_MPC").resolve()
    if not day_out.is_dir():
        raise FileNotFoundError(f"Day output directory not found: {day_out}")
    df = load_manifest(day_out)

    prs = Presentation()
    add_title_slide(prs, args.title, args.subtitle, date_label)

    layout_cfg = dict(
        left_margin_in=float(args.left_margin_in),
        top_margin_in=float(args.top_margin_in),
        table_width_in=float(args.table_width_in),
        table_height_in=float(args.table_height_in),
        gif_width_in=float(args.gif_width_in),
        gif_gap_in=float(args.gif_gap_in),
        max_gifs_per_slide=int(args.max_gifs_per_slide),
        rows_per_table=int(args.rows_per_table),
        day_out=day_out,
    )

    # Determine requested columns: CLI has priority; otherwise use NEEDED
    requested_cols = args.cols if args.cols else NEEDED

    # Group by (object, site)
    grouped = []
    for (obj, site), g in df.groupby(["object", "site"], sort=True):
        grouped.append(((obj, site), g.sort_values("sub_id")))
    if args.limit_groups is not None:
        grouped = grouped[: int(args.limit_groups)]

    for key, gdf in grouped:
        add_group_slide(prs, key, gdf, requested_cols, layout_cfg)

    # Output path
    ppt_outdir = Path(args.ppt_outdir).resolve() if args.ppt_outdir else thumbs_root
    ppt_outdir.mkdir(parents=True, exist_ok=True)
    ppt_name = args.ppt_name_template.format(date=date_label)
    out_path = ppt_outdir / ppt_name
    prs.save(str(out_path))
    print(f"  Saved PPT: {out_path}")
    print(f"  Source day dir: {day_out}")
    print(f"  Groups (object, site): {len(grouped)}")


if __name__ == "__main__":
    main()
